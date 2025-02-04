# from django.shortcuts import render, redirect, get_object_or_404
# from django.http import HttpResponse, FileResponse
# from django.core.mail import send_mail
# from django.conf import settings
# from django.contrib import messages
# from django.contrib.auth.decorators import login_required
# from .models import Event, Participant
# import pandas as pd
# from .convter import ppt2pdf
# from pptx import Presentation
# from django.core.mail import send_mail, EmailMessage
# import requests
# import os

# def index(request):
#     return render(request, 'index.html')

# @login_required
# def create(request):
#     if request.method == "POST":
#         csv = request.FILES.get('csv')
#         temp = request.FILES.get('template')
#         event_name = request.POST.get('event_name')
        
#         event = Event(user=request.user,
#                       event_name=event_name,
#                       csv_file=csv,
#                       template=temp)
#         event.save()

#         return redirect(f"/certificate/{event.id}/{event.slug}")

#     return render(request, 'certificate/create_event.html')

# @login_required
# def delete_event(request, id, slug):
#     event = Event.objects.filter(slug=slug, id=id).first()
#     if event.user == request.user:
#         event.delete()
#     return redirect('view_certificate_status')

# @login_required
# def track(request, id, slug):
#     event = Event.objects.filter(slug=slug, id=id).first()
#     if event.message:
#         return render(request, 'certificate/track.html', {
#             'event_name': event.event_name,
#             'event_date': event.date,
#             'participat_details': Participant.objects.filter(event=event)
#         })

#     prs = Presentation(event.template)
#     st = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 st = st + shape.text
#                 st = st + " "
#     li = st.split()
#     tags = []
#     for i in li:
#         if i[0] == "<" and i[-1] == ">":
#             tags.append(i)

#     if request.method == "POST":
#         email_col = request.POST.get('emails')
#         subject = request.POST.get('subject')
#         mess = request.POST.get('mess')
#         values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]    
        
#         event.email_column = email_col
#         event.message = mess
#         event.subject = subject
#         event.save()

#         df = pd.read_csv(event.csv_file)
#         df_len = df.shape[0]
#         i = 0

#         data = {
#             "client_id": "473692605888-cju6d348n671vao051da72ct265htjkb.apps.googleusercontent.com",
#             "client_secret": "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",
#             "refresh_token": "1//04T-bb2HQppWmCgYIARAAGAQSNwF-L9IrY7kPEd1yOismONmGZO4VTfNaZbAqsvnop33PqxDbCxquZutR4803mOGoBSOY6BaZcM8",
#             'grant_type': 'refresh_token'
#         }
#         a = requests.post("https://www.googleapis.com/oauth2/v4/token", data)
#         token = f"Bearer {dict(a.json()).get('access_token')}"

#         root_directory = settings.BASE_DIR  # Change this to your desired root directory
        
#         while i < df_len:
#             prs = Presentation(event.template)
#             j = ""
#             if i < 9:
#                 j = "00"
#             elif i >= 9 and i < 99:
#                 j = "0"
            
#             for tag, v_type, value in values:
#                 for slide in prs.slides:
#                     for shape in slide.shapes:
#                         if shape.has_text_frame:
#                             if(shape.text.find(tag)) != -1:
#                                     text_frame = shape.text_frame
#                                     for paragraph in text_frame.paragraphs:
#                                         for run in paragraph.runs:
#                                             cur_text = run.text
#                                             if v_type == 'text':
#                                                 new_text = cur_text.replace(tag, value)
#                                             elif v_type == 'date':
#                                                 new_text = cur_text.replace(tag, '/'.join(value.split('-')[::-1]))
#                                             elif v_type == 'csv':
#                                                 new_text = cur_text.replace(tag, df.loc[i, value])
#                                             elif v_type == "auto":
#                                                 new_text = cur_text.replace(tag, value+"/"+j+str(i+1))
#                                             else:
#                                                 pass
#                                             run.text = new_text
                                            
#             s_name = df.loc[i, event.email_column].split('@')[0]
#             prs.save(os.path.join(root_directory, f"{s_name}.pptx"))
#             f_id = ppt2pdf(os.path.join(root_directory, f"{s_name}.pptx"), s_name, token)
#             download_url = f"https://docs.google.com/presentation/d/{f_id}/export/pdf"

#             r = requests.get(download_url, allow_redirects=True)
#             certificate_path = os.path.join(root_directory, f"{s_name}_certificate.pdf")
#             open(certificate_path, 'wb').write(r.content)

#             try:
#                 mail = EmailMessage(subject,
#                                     f"Hello, {s_name} \n{mess}",
#                                     settings.EMAIL_HOST_USER,
#                                     [df.loc[i, event.email_column]])
#                 mail.attach_file(certificate_path)
#                 mail.send()
#                 Participant(event=event, email=df.loc[i, event.email_column], status=True).save()

#                 # Prepare response to download certificate
#                 certificate_name = f"{s_name}_certificate.pdf"
#                 response = FileResponse(open(certificate_path, 'rb'), content_type='application/pdf')
#                 response['Content-Disposition'] = f'attachment; filename="{certificate_name}"'

#                 # Remove certificate file after sending email
#                 os.remove(certificate_path)
#                 os.remove(os.path.join(root_directory, f"{s_name}.pptx"))

#                 return response  # Return response for download
#             except:
#                 Participant(event=event, email=df.loc[i, event.email_column], status=False).save()

#             i = i + 1

#         messages.success(request, "Certificates Sent Successfully !!")
#         return redirect(f"/certificate/{event.id}/{event.slug}")

#     return render(request, 'certificate/map_tags_of_template.html', {
#         'columns': list(pd.read_csv(event.csv_file).columns),
#         'tags': tags,
#     })

# @login_required
# def view_certificate_status(request):
#     return render(request, 'certificate/view_certificate_status.html', {
#         'events': Event.objects.filter(user=request.user)
#     })











# from django.shortcuts import render, redirect, get_object_or_404
# from django.http import HttpResponse
# from django.core.mail import send_mail
# from django.conf import settings
# from django.contrib import messages
# from django.contrib.auth.decorators import login_required
# from .models import Event, Participant
# import pandas as pd
# from .convter import ppt2pdf
# from pptx import Presentation
# from django.core.mail import send_mail, EmailMessage
# import requests
# import os

# def index(request):
# 	return render(request, 'index.html')

# @login_required
# def create(request):
# 	if request.method == "POST":
# 		csv = request.FILES.get('csv')
# 		temp = request.FILES.get('template')
# 		event_name = request.POST.get('event_name')
		
# 		event = Event(user = request.user,
# 			event_name = event_name,
# 			csv_file = csv,
# 			template = temp)
# 		event.save()

# 		return redirect(f"/certificate/{event.id}/{event.slug}")

# 	return render(request, 'certificate/create_event.html')

# @login_required
# def delete_event(request, id, slug):
# 	event = Event.objects.filter(slug=slug, id=id).first()
# 	if event.user == request.user:
# 	    event.delete()
# 	return redirect('view_certificate_status')

# @login_required
# def track(request, id, slug):
# 	event = Event.objects.filter(slug=slug, id=id).first()
# 	if event.message:

# 		return render(request, 'certificate/track.html', {
# 			'event_name': event.event_name,
# 			'event_date': event.date,
# 			'participat_details': Participant.objects.filter(event=event)
# 			})

# 	prs = Presentation(event.template)
# 	st=""
# 	for slide in prs.slides:
# 		for shape in slide.shapes:
# 			if shape.has_text_frame:
# 				st = st + shape.text
# 				st = st + " "
# 	li = st.split()
# 	tags = []
# 	for i in li:
# 		if i[0] == "<" and i[-1] == ">":
# 			tags.append(i)

# 	if request.method == "POST":
# 		email_col = request.POST.get('emails')
# 		subject = request.POST.get('subject')
# 		mess = request.POST.get('mess')
# 		values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]	
		
# 		event.email_column = email_col
# 		event.message = mess
# 		event.subject = subject
# 		event.save()

# 		df=pd.read_csv(event.csv_file)
# 		df_len=df.shape
# 		i=0

# 		data = {
#         	"client_id":"YOUR_GOOGLE_API_CLIENT_ID",
#         	"client_secret":"YOUR_CLIENT_SECRET",
#         	"refresh_token": "GOOGLE DRIVE REFRESH TOKEN",
#         	'grant_type': 'refresh_token'
#             }
# 		a = requests.post("https://www.googleapis.com/oauth2/v4/token", data)
# 		token = f"Bearer {dict(a.json()).get('access_token')}"
# 		li=["First","Second","Third"]
# 		while i < df_len[0]:
# 			prs = Presentation(event.template)
# 			j=""
# 			if i<9:
# 				j="00"
# 			elif i>=9 and i < 99 :
# 				j="0"
			
# 			for tag, v_type, value in values:
# 				for slide in prs.slides:
# 					for shape in slide.shapes:
# 						if shape.has_text_frame:
# 							if(shape.text.find(tag))!=-1:
# 									text_frame = shape.text_frame
# 									for paragraph in text_frame.paragraphs:
# 										for run in paragraph.runs:
# 											cur_text = run.text
# 											if v_type == 'text':
# 												new_text = cur_text.replace(tag, value)
# 											elif v_type == 'date':
# 												new_text = cur_text.replace(tag, '/'.join(value.split('-')[::-1]))
# 											elif v_type == 'csv':
# 												new_text = cur_text.replace(tag, df.loc[i,value])
# 											elif v_type == "auto":
# 												new_text = cur_text.replace(tag, value+"/"+j+str(i+1))
# 											else:
# 												pass
# 											run.text = new_text
											
			
# 			s_name = df.loc[i,event.email_column].split('@')[0]
# 			prs.save(s_name+".pptx")
# 			f_id = ppt2pdf(s_name+".pptx",s_name, token)
# 			r = requests.get(f"https://docs.google.com/presentation/d/{f_id}/export/pdf", allow_redirects=True)
# 			open(s_name+'.pdf', 'wb').write(r.content)

# 			try:
# 				mail = EmailMessage(subject,
# 					f"Hello, {s_name} \n{mess}",
# 					settings.EMAIL_HOST_USER,
# 					[df.loc[i,event.email_column]])
# 				mail.attach_file(s_name+'.pdf')
# 				mail.send()
# 				Participant(event=event, email=df.loc[i,event.email_column], status=True).save()
# 				os.remove(s_name+'.pdf')
# 				os.remove(s_name+".pptx")
# 			except:
# 				Participant(event=event, email=df.loc[i,event.email_column], status=False).save()
# 				os.remove(s_name+'.pdf')
# 				os.remove(s_name+".pptx")
# 			i=i+1

# 		messages.success(request, "Certificates Sent Successfuly !!")
# 		return redirect(f"/certificate/{event.id}/{event.slug}")


# 	return render(request, 'certificate/map_tags_of_template.html',{
# 		'columns': list(pd.read_csv(event.csv_file).columns),
# 		'tags': tags,
# 		})


# @login_required
# def view_certificate_status(request):
# 	return render(request, 'certificate/view_certificate_status.html',{
# 		'events': Event.objects.filter(user=request.user)
# 		})



# from django.shortcuts import render, redirect, get_object_or_404
# from django.http import HttpResponse, FileResponse
# from django.core.mail import send_mail
# from django.conf import settings
# from django.contrib import messages
# from django.contrib.auth.decorators import login_required
# from .models import Event, Participant
# import pandas as pd
# from .convter import ppt2pdf
# from pptx import Presentation
# from django.core.mail import send_mail, EmailMessage
# import requests
# import os

# def index(request):
#     return render(request, 'index.html')

# @login_required
# def create(request):
#     if request.method == "POST":
#         csv = request.FILES.get('csv')
#         temp = request.FILES.get('template')
#         event_name = request.POST.get('event_name')
        
#         event = Event(user=request.user,
#                       event_name=event_name,
#                       csv_file=csv,
#                       template=temp)
#         event.save()

#         return redirect(f"/certificate/{event.id}/{event.slug}")

#     return render(request, 'certificate/create_event.html')

# @login_required
# def delete_event(request, id, slug):
#     event = Event.objects.filter(slug=slug, id=id).first()
#     if event.user == request.user:
#         event.delete()
#     return redirect('view_certificate_status')

# @login_required
# def track(request, id, slug):
#     event = Event.objects.filter(slug=slug, id=id).first()
#     if event.message:
#         return render(request, 'certificate/track.html', {
#             'event_name': event.event_name,
#             'event_date': event.date,
#             'participat_details': Participant.objects.filter(event=event)
#         })

#     prs = Presentation(event.template)
#     st = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 st = st + shape.text
#                 st = st + " "
#     li = st.split()
#     tags = []
#     for i in li:
#         if i[0] == "<" and i[-1] == ">":
#             tags.append(i)

#     if request.method == "POST":
#         email_col = request.POST.get('emails')
#         subject = request.POST.get('subject')
#         mess = request.POST.get('mess')
#         values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]    
        
#         event.email_column = email_col
#         event.message = mess
#         event.subject = subject
#         event.save()

#         df = pd.read_csv(event.csv_file)
#         df_len = df.shape[0]
#         i = 0

#         # credetials for tdata0465@gmail.com 


#         data = {
#             "client_id": "473692605888-cju6d348n671vao051da72ct265htjkb.apps.googleusercontent.com",
#             "client_secret": "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",
#             "refresh_token": "1//04T-bb2HQppWmCgYIARAAGAQSNwF-L9IrY7kPEd1yOismONmGZO4VTfNaZbAqsvnop33PqxDbCxquZutR4803mOGoBSOY6BaZcM8",
#             'grant_type': 'refresh_token'
#         }
#         a = requests.post("https://www.googleapis.com/oauth2/v4/token", data)
#         token = f"Bearer {dict(a.json()).get('access_token')}"

#         root_directory = settings.BASE_DIR  # Change this to your desired root directory
        
#         while i < df_len:
#             prs = Presentation(event.template)
#             j = ""
#             if i < 9:
#                 j = "00"
#             elif i >= 9 and i < 99:
#                 j = "0"
            
#             for tag, v_type, value in values:
#                 for slide in prs.slides:
#                     for shape in slide.shapes:
#                         if shape.has_text_frame:
#                             if(shape.text.find(tag)) != -1:
#                                     text_frame = shape.text_frame
#                                     for paragraph in text_frame.paragraphs:
#                                         for run in paragraph.runs:
#                                             cur_text = run.text
#                                             if v_type == 'text':
#                                                 new_text = cur_text.replace(tag, value)
#                                             elif v_type == 'date':
#                                                 new_text = cur_text.replace(tag, '/'.join(value.split('-')[::-1]))
#                                             elif v_type == 'csv':
#                                                 new_text = cur_text.replace(tag, df.loc[i, value])
#                                             elif v_type == "auto":
#                                                 new_text = cur_text.replace(tag, value+"/"+j+str(i+1))
#                                             else:
#                                                 pass
#                                             run.text = new_text
                                            
#             s_name = df.loc[i, event.email_column].split('@')[0]
#             prs.save(os.path.join(root_directory, f"{s_name}.pptx"))
#             f_id = ppt2pdf(os.path.join(root_directory, f"{s_name}.pptx"), s_name, token)
#             download_url = f"https://docs.google.com/presentation/d/{f_id}/export/pdf"

#             r = requests.get(download_url, allow_redirects=True)
#             certificate_path = os.path.join(root_directory, f"{s_name}_certificate.pdf")
#             open(certificate_path, 'wb').write(r.content)

#             try:
#                 mail = EmailMessage(subject,
#                                     f"Hello, {s_name} \n{mess}\n\nThis service is powered by EDEMS PVT Ltd.",
#                                     settings.EMAIL_HOST_USER,
#                                     [df.loc[i, event.email_column]])
#                 mail.attach_file(certificate_path)
#                 mail.send()
#                 Participant(event=event, email=df.loc[i, event.email_column], status=True).save()

#                 # Prepare response to download certificate
#                 certificate_name = f"{s_name}_certificate.pdf"
#                 response = FileResponse(open(certificate_path, 'rb'), content_type='application/pdf')
#                 response['Content-Disposition'] = f'attachment; filename="{certificate_name}"'

#                 # Remove certificate file after sending email
#                 os.remove(certificate_path)
#                 os.remove(os.path.join(root_directory, f"{s_name}.pptx"))

#                 return response  # Return response for download
#             except:
#                 Participant(event=event, email=df.loc[i, event.email_column], status=False).save()

#             i = i + 1

#         messages.success(request, "Certificates Sent Successfully !!")
#         return redirect(f"/certificate/{event.id}/{event.slug}")

#     return render(request, 'certificate/map_tags_of_template.html', {
#         'columns': list(pd.read_csv(event.csv_file).columns),
#         'tags': tags,
#     })


# @login_required
# def view_certificate_status(request):
#     return render(request, 'certificate/view_certificate_status.html', {
#         'events': Event.objects.filter(user=request.user)
#     })




# image processing logic 




# from django.shortcuts import render, redirect, get_object_or_404
# from django.http import HttpResponse, FileResponse
# from django.core.mail import EmailMessage
# from django.conf import settings
# from django.contrib import messages
# from django.contrib.auth.decorators import login_required
# from .models import Event, Participant
# import pandas as pd
# from .convter import ppt2pdf
# from pptx import Presentation
# import requests
# import os
# import urllib.request
# from PIL import Image
# from io import BytesIO

# def index(request):
#     return render(request, 'index.html')

# @login_required
# def create(request):
#     if request.method == "POST":
#         csv = request.FILES.get('csv')
#         temp = request.FILES.get('template')
#         event_name = request.POST.get('event_name')
        
#         event = Event(user=request.user,
#                       event_name=event_name,
#                       csv_file=csv,
#                       template=temp)
#         event.save()

#         return redirect(f"/certificate/{event.id}/{event.slug}")

#     return render(request, 'certificate/create_event.html')

# @login_required
# def delete_event(request, id, slug):
#     event = Event.objects.filter(slug=slug, id=id).first()
#     if event.user == request.user:
#         event.delete()
#     return redirect('view_certificate_status')

# @login_required
# def track(request, id, slug):
#     event = Event.objects.filter(slug=slug, id=id).first()
#     if event.message:
#         return render(request, 'certificate/track.html', {
#             'event_name': event.event_name,
#             'event_date': event.date,
#             'participat_details': Participant.objects.filter(event=event)
#         })

#     prs = Presentation(event.template)
#     st = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 st += shape.text + " "
    
#     li = st.split()
#     tags = [i for i in li if i[0] == "<" and i[-1] == ">"]

#     if request.method == "POST":
#         email_col = request.POST.get('emails')
#         subject = request.POST.get('subject')
#         mess = request.POST.get('mess')
#         values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]
        
#         event.email_column = email_col
#         event.message = mess
#         event.subject = subject
#         event.save()

#         df = pd.read_csv(event.csv_file)
#         df_len = df.shape[0]
#         i = 0

#         data = {
#             "client_id": "473692605888-cju6d348n671vao051da72ct265htjkb.apps.googleusercontent.com",
#             "client_secret": "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",
#             "refresh_token": "1//04T-bb2HQppWmCgYIARAAGAQSNwF-L9IrY7kPEd1yOismONmGZO4VTfNaZbAqsvnop33PqxDbCxquZutR4803mOGoBSOY6BaZcM8",
#             'grant_type': 'refresh_token'
#         }
#         a = requests.post("https://www.googleapis.com/oauth2/v4/token", data)
#         token = f"Bearer {dict(a.json()).get('access_token')}"

#         root_directory = settings.BASE_DIR
        
#         while i < df_len:
#             prs = Presentation(event.template)
#             j = "00" if i < 9 else "0" if i < 99 else ""

#             for tag, v_type, value in values:
#                 for slide in prs.slides:
#                     for shape in slide.shapes:
#                         if shape.has_text_frame:
#                             if tag in shape.text:
#                                 text_frame = shape.text_frame
#                                 for paragraph in text_frame.paragraphs:
#                                     for run in paragraph.runs:
#                                         cur_text = run.text
#                                         if v_type == 'text':
#                                             new_text = cur_text.replace(tag, value)
#                                         elif v_type == 'date':
#                                             new_text = cur_text.replace(tag, '/'.join(value.split('-')[::-1]))
#                                         elif v_type == 'csv':
#                                             new_text = cur_text.replace(tag, df.loc[i, value])
#                                         elif v_type == "auto":
#                                             new_text = cur_text.replace(tag, value + "/" + j + str(i + 1))
#                                         elif v_type == "image":  # Handling image URLs
#                                             img_url = df.loc[i, value]
#                                             urllib.request.urlretrieve(img_url, "temp_image.png")
#                                             left = Inches(1)  # Adjust position
#                                             top = Inches(1)   # Adjust position
#                                             slide.shapes.add_picture("temp_image.png", left, top, width=Inches(2), height=Inches(2))
#                                             os.remove("temp_image.png")  # Clean up temp image
#                                         else:
#                                             continue
#                                         run.text = new_text
            
#             s_name = df.loc[i, event.email_column].split('@')[0]
#             prs.save(os.path.join(root_directory, f"{s_name}.pptx"))
#             f_id = ppt2pdf(os.path.join(root_directory, f"{s_name}.pptx"), s_name, token)
#             download_url = f"https://docs.google.com/presentation/d/{f_id}/export/pdf"

#             r = requests.get(download_url, allow_redirects=True)
#             certificate_path = os.path.join(root_directory, f"{s_name}_certificate.pdf")
#             open(certificate_path, 'wb').write(r.content)

#             try:
#                 mail = EmailMessage(subject,
#                                     f"Hello, {s_name} \n{mess}\n\nThis service is powered by EDEMS PVT Ltd.",
#                                     settings.EMAIL_HOST_USER,
#                                     [df.loc[i, event.email_column]])
#                 mail.attach_file(certificate_path)
#                 mail.send()
#                 Participant(event=event, email=df.loc[i, event.email_column], status=True).save()

#                 # Prepare response to download certificate
#                 response = FileResponse(open(certificate_path, 'rb'), content_type='application/pdf')
#                 response['Content-Disposition'] = f'attachment; filename="{s_name}_certificate.pdf"'

#                 # Remove certificate file after sending email
#                 os.remove(certificate_path)
#                 os.remove(os.path.join(root_directory, f"{s_name}.pptx"))

#                 return response  # Return response for download
#             except Exception as e:
#                 Participant(event=event, email=df.loc[i, event.email_column], status=False).save()

#             i += 1

#         messages.success(request, "Certificates Sent Successfully !!")
#         return redirect(f"/certificate/{event.id}/{event.slug}")

#     return render(request, 'certificate/map_tags_of_template.html', {
#         'columns': list(pd.read_csv(event.csv_file).columns),
#         'tags': tags,
#     })

# @login_required
# def view_certificate_status(request):
#     return render(request, 'certificate/view_certificate_status.html', {
#         'events': Event.objects.filter(user=request.user)
#     })




# from django.shortcuts import render, redirect, get_object_or_404
# from django.http import HttpResponse, FileResponse
# from .converter import ppt2pdf

# from django.core.mail import send_mail
# from django.conf import settings
# from django.contrib import messages
# from django.contrib.auth.decorators import login_required
# from .models import Event, Participant
# import pandas as pd
# from .converter import ppt2pdf
# from pptx import Presentation
# from django.core.mail import send_mail, EmailMessage
# import requests
# import os
# import qrcode
# from PIL import Image

# def index(request):
#     return render(request, 'index.html')

# @login_required
# def create(request):
#     if request.method == "POST":
#         csv = request.FILES.get('csv')
#         temp = request.FILES.get('template')
#         event_name = request.POST.get('event_name')
        
#         event = Event(user=request.user,
#                       event_name=event_name,
#                       csv_file=csv,
#                       template=temp)
#         event.save()

#         return redirect(f"/certificate/{event.id}/{event.slug}")

#     return render(request, 'certificate/create_event.html')

# @login_required
# def delete_event(request, id, slug):
#     event = Event.objects.filter(slug=slug, id=id).first()
#     if event.user == request.user:
#         event.delete()
#     return redirect('view_certificate_status')

# @login_required
# def track(request, id, slug):
#     event = Event.objects.filter(slug=slug, id=id).first()
#     if event.message:
#         return render(request, 'certificate/track.html', {
#             'event_name': event.event_name,
#             'event_date': event.date,
#             'participat_details': Participant.objects.filter(event=event)
#         })

#     prs = Presentation(event.template)
#     st = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 st = st + shape.text
#                 st = st + " "
#     li = st.split()
#     tags = []
#     for i in li:
#         if i[0] == "<" and i[-1] == ">":
#             tags.append(i)

#     if request.method == "POST":
#         email_col = request.POST.get('emails')
#         subject = request.POST.get('subject')
#         mess = request.POST.get('mess')
#         values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]    
        
#         event.email_column = email_col
#         event.message = mess
#         event.subject = subject
#         event.save()

#         df = pd.read_csv(event.csv_file)
#         df_len = df.shape[0]
#         i = 0

#         data = {
#             "client_id": "473692605888-cju6d348n671vao051da72ct265htjkb.apps.googleusercontent.com",
#             "client_secret": "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",
#             "refresh_token": "1//04T-bb2HQppWmCgYIARAAGAQSNwF-L9IrY7kPEd1yOismONmGZO4VTfNaZbAqsvnop33PqxDbCxquZutR4803mOGoBSOY6BaZcM8",
#             'grant_type': 'refresh_token'
#         }
#         a = requests.post("https://www.googleapis.com/oauth2/v4/token", data)
#         token = f"Bearer {dict(a.json()).get('access_token')}"

#         root_directory = settings.BASE_DIR  # Change this to your desired root directory
        
#         while i < df_len:
#             prs = Presentation(event.template)
#             j = ""
#             if i < 9:
#                 j = "00"
#             elif i >= 9 and i < 99:
#                 j = "0"
            
#             for tag, v_type, value in values:
#                 for slide in prs.slides:
#                     for shape in slide.shapes:
#                         if shape.has_text_frame:
#                             if(shape.text.find(tag)) != -1:
#                                 text_frame = shape.text_frame
#                                 for paragraph in text_frame.paragraphs:
#                                     for run in paragraph.runs:
#                                         cur_text = run.text
#                                         if v_type == 'text':
#                                             new_text = cur_text.replace(tag, value)
#                                         elif v_type == 'date':
#                                             new_text = cur_text.replace(tag, '/'.join(value.split('-')[::-1]))
#                                         elif v_type == 'csv':
#                                             new_text = cur_text.replace(tag, df.loc[i, value])
#                                         elif v_type == "auto":
#                                             new_text = cur_text.replace(tag, value+"/"+j+str(i+1))
#                                         elif v_type == "qr":
#                                             # Generate QR code
#                                             qr_data = df.loc[i, value]
#                                             qr_img = qrcode.make(qr_data)
#                                             qr_img_path = os.path.join(root_directory, f"{value}_{i}.png")
#                                             qr_img.save(qr_img_path)

#                                             # Add QR code to the slide
#                                             left = top = width = height = 0 # Set appropriate positions and size
#                                             slide.shapes.add_picture(qr_img_path, left, top, width, height)
#                                             os.remove(qr_img_path)  # Remove the QR image after use
#                                         else:
#                                             pass
#                                         run.text = new_text
                                        
#             s_name = df.loc[i, event.email_column].split('@')[0]
#             prs.save(os.path.join(root_directory, f"{s_name}.pptx"))
#             f_id = ppt2pdf(os.path.join(root_directory, f"{s_name}.pptx"), s_name, token)
#             download_url = f"https://docs.google.com/presentation/d/{f_id}/export/pdf"

#             r = requests.get(download_url, allow_redirects=True)
#             certificate_path = os.path.join(root_directory, f"{s_name}_certificate.pdf")
#             open(certificate_path, 'wb').write(r.content)

#             try:
#                 mail = EmailMessage(subject,
#                                     f"Hello, {s_name} \n{mess}\n\nThis service is powered by EDEMS PVT Ltd.",
#                                     settings.EMAIL_HOST_USER,
#                                     [df.loc[i, event.email_column]])
#                 mail.attach_file(certificate_path)
#                 mail.send()
#                 Participant(event=event, email=df.loc[i, event.email_column], status=True).save()

#                 # Prepare response to download certificate
#                 certificate_name = f"{s_name}_certificate.pdf"
#                 response = FileResponse(open(certificate_path, 'rb'), content_type='application/pdf')
#                 response['Content-Disposition'] = f'attachment; filename="{certificate_name}"'

#                 # Remove certificate file after sending email
#                 os.remove(certificate_path)
#                 os.remove(os.path.join(root_directory, f"{s_name}.pptx"))

#                 return response  # Return response for download
#             except:
#                 Participant(event=event, email=df.loc[i, event.email_column], status=False).save()

#             i = i + 1

#         messages.success(request, "Certificates Sent Successfully !!")
#         return redirect(f"/certificate/{event.id}/{event.slug}")

#     return render(request, 'certificate/map_tags_of_template.html', {
#         'columns': list(pd.read_csv(event.csv_file).columns),
#         'tags': tags,
#     })

# @login_required
# def view_certificate_status(request):
#     return render(request, 'certificate/view_certificate_status.html', {
#         'events': Event.objects.filter(user=request.user)
#     })





# from django.shortcuts import render, redirect, get_object_or_404
# from django.http import HttpResponse, FileResponse
# from django.core.mail import send_mail
# from django.conf import settings
# from django.contrib import messages
# from django.contrib.auth.decorators import login_required
# from .models import Event, Participant
# import pandas as pd
# from .convter import ppt2pdf
# from pptx import Presentation
# from django.core.mail import send_mail, EmailMessage
# import requests
# import os

# def index(request):
#     return render(request, 'index.html')

# @login_required
# def create(request):
#     if request.method == "POST":
#         csv = request.FILES.get('csv')
#         temp = request.FILES.get('template')
#         event_name = request.POST.get('event_name')
        
#         event = Event(user=request.user,
#                       event_name=event_name,
#                       csv_file=csv,
#                       template=temp)
#         event.save()

#         return redirect(f"/certificate/{event.id}/{event.slug}")

#     return render(request, 'certificate/create_event.html')

# @login_required
# def delete_event(request, id, slug):
#     event = Event.objects.filter(slug=slug, id=id).first()
#     if event.user == request.user:
#         event.delete()
#     return redirect('view_certificate_status')

# @login_required
# def track(request, id, slug):
#     event = Event.objects.filter(slug=slug, id=id).first()
#     if event.message:
#         return render(request, 'certificate/track.html', {
#             'event_name': event.event_name,
#             'event_date': event.date,
#             'participat_details': Participant.objects.filter(event=event)
#         })

#     prs = Presentation(event.template)
#     st = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 st = st + shape.text
#                 st = st + " "
#     li = st.split()
#     tags = []
#     for i in li:
#         if i[0] == "<" and i[-1] == ">":
#             tags.append(i)

#     if request.method == "POST":
#         email_col = request.POST.get('emails')
#         subject = request.POST.get('subject')
#         mess = request.POST.get('mess')
#         values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]    
        
#         event.email_column = email_col
#         event.message = mess
#         event.subject = subject
#         event.save()

#         df = pd.read_csv(event.csv_file)
#         df_len = df.shape[0]
#         i = 0

#         data = {
#             "client_id": "473692605888-cju6d348n671vao051da72ct265htjkb.apps.googleusercontent.com",
#             "client_secret": "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",
#             "refresh_token": "1//046FabrbGAKw5CgYIARAAGAQSNwF-L9Ir_VoquOar3pXx8ZlrCsD0WRZVD7spGGBU6GULCNzt9DQkjPCnE5fI_fs05W6HTlMkhSw",
#             'grant_type': 'refresh_token'
#         }
#         a = requests.post("https://www.googleapis.com/oauth2/v4/token", data)
#         token = f"Bearer {dict(a.json()).get('access_token')}"

#         root_directory = settings.BASE_DIR  # Change this to your desired root directory
        
#         while i < df_len:
#             prs = Presentation(event.template)
#             j = ""
#             if i < 9:
#                 j = "00"
#             elif i >= 9 and i < 99:
#                 j = "0"
            
#             for tag, v_type, value in values:
#                 for slide in prs.slides:
#                     for shape in slide.shapes:
#                         if shape.has_text_frame:
#                             if(shape.text.find(tag)) != -1:
#                                     text_frame = shape.text_frame
#                                     for paragraph in text_frame.paragraphs:
#                                         for run in paragraph.runs:
#                                             cur_text = run.text
#                                             if v_type == 'text':
#                                                 new_text = cur_text.replace(tag, value)
#                                             elif v_type == 'date':
#                                                 new_text = cur_text.replace(tag, '/'.join(value.split('-')[::-1]))
#                                             elif v_type == 'csv':
#                                                 new_text = cur_text.replace(tag, df.loc[i, value])
#                                             elif v_type == "auto":
#                                                 new_text = cur_text.replace(tag, value+"/"+j+str(i+1))
#                                             else:
#                                                 pass
#                                             run.text = new_text
                                            
#             s_name = df.loc[i, event.email_column].split('@')[0]
#             prs.save(os.path.join(root_directory, f"{s_name}.pptx"))
#             f_id = ppt2pdf(os.path.join(root_directory, f"{s_name}.pptx"), s_name, token)
#             download_url = f"https://docs.google.com/presentation/d/{f_id}/export/pdf"

#             r = requests.get(download_url, allow_redirects=True)
#             certificate_path = os.path.join(root_directory, f"{s_name}_certificate.pdf")
#             open(certificate_path, 'wb').write(r.content)

#             try:
#                 mail = EmailMessage(subject,
#                                     f"Hello, {s_name} \n{mess}",
#                                     settings.EMAIL_HOST_USER,
#                                     [df.loc[i, event.email_column]])
#                 mail.attach_file(certificate_path)
#                 mail.send()
#                 Participant(event=event, email=df.loc[i, event.email_column], status=True).save()

#                 # Prepare response to download certificate
#                 certificate_name = f"{s_name}_certificate.pdf"
#                 response = FileResponse(open(certificate_path, 'rb'), content_type='application/pdf')
#                 response['Content-Disposition'] = f'attachment; filename="{certificate_name}"'

#                 # Remove certificate file after sending email
#                 os.remove(certificate_path)
#                 os.remove(os.path.join(root_directory, f"{s_name}.pptx"))

#                 return response  # Return response for download
#             except:
#                 Participant(event=event, email=df.loc[i, event.email_column], status=False).save()

#             i = i + 1

#         messages.success(request, "Certificates Sent Successfully !!")
#         return redirect(f"/certificate/{event.id}/{event.slug}")

#     return render(request, 'certificate/map_tags_of_template.html', {
#         'columns': list(pd.read_csv(event.csv_file).columns),
#         'tags': tags,
#     })

# @login_required
# def view_certificate_status(request):
#     return render(request, 'certificate/view_certificate_status.html', {
#         'events': Event.objects.filter(user=request.user)
#     })






from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse, FileResponse
from django.core.mail import EmailMessage
from django.conf import settings
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from .models import Event, Participant
import pandas as pd
from .convter import ppt2pdf
from pptx import Presentation
import requests
import os
import urllib.request
from PIL import Image
from io import BytesIO

def index(request):
    return render(request, 'index.html')

@login_required
def create(request):
    if request.method == "POST":
        csv = request.FILES.get('csv')
        temp = request.FILES.get('template')
        event_name = request.POST.get('event_name')
        
        event = Event(user=request.user,
                      event_name=event_name,
                      csv_file=csv,
                      template=temp)
        event.save()

        return redirect(f"/certificate/{event.id}/{event.slug}")

    return render(request, 'certificate/create_event.html')

@login_required
def delete_event(request, id, slug):
    event = Event.objects.filter(slug=slug, id=id).first()
    if event.user == request.user:
        event.delete()
    return redirect('view_certificate_status')

@login_required
def track(request, id, slug):
    event = Event.objects.filter(slug=slug, id=id).first()
    if event.message:
        return render(request, 'certificate/track.html', {
            'event_name': event.event_name,
            'event_date': event.date,
            'participat_details': Participant.objects.filter(event=event)
        })

    prs = Presentation(event.template)
    st = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                st += shape.text + " "
    
    li = st.split()
    tags = [i for i in li if i[0] == "<" and i[-1] == ">"]

    if request.method == "POST":
        email_col = request.POST.get('emails')
        subject = request.POST.get('subject')
        mess = request.POST.get('mess')
        values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]
        
        event.email_column = email_col
        event.message = mess
        event.subject = subject
        event.save()

        df = pd.read_csv(event.csv_file)
        df_len = df.shape[0]
        i = 0

        data = {
            "client_id": "473692605888-cju6d348n671vao051da72ct265htjkb.apps.googleusercontent.com",
            "client_secret": "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",
            "refresh_token": "1//046FabrbGAKw5CgYIARAAGAQSNwF-L9Ir_VoquOar3pXx8ZlrCsD0WRZVD7spGGBU6GULCNzt9DQkjPCnE5fI_fs05W6HTlMkhSw",
            'grant_type': 'refresh_token'
        }
        a = requests.post("https://www.googleapis.com/oauth2/v4/token", data)
        token = f"Bearer {dict(a.json()).get('access_token')}"

        root_directory = settings.BASE_DIR
        
        while i < df_len:
            prs = Presentation(event.template)
            j = "00" if i < 9 else "0" if i < 99 else ""

            for tag, v_type, value in values:
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            if tag in shape.text:
                                text_frame = shape.text_frame
                                for paragraph in text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        cur_text = run.text
                                        if v_type == 'text':
                                            new_text = cur_text.replace(tag, value)
                                        elif v_type == 'date':
                                            new_text = cur_text.replace(tag, '/'.join(value.split('-')[::-1]))
                                        elif v_type == 'csv':
                                            new_text = cur_text.replace(tag, df.loc[i, value])
                                        elif v_type == "auto":
                                            new_text = cur_text.replace(tag, value + "/" + j + str(i + 1))
                                        elif v_type == "image":  # Handling image URLs
                                            img_url = df.loc[i, value]
                                            urllib.request.urlretrieve(img_url, "temp_image.png")
                                            left = Inches(1)  # Adjust position
                                            top = Inches(1)   # Adjust position
                                            slide.shapes.add_picture("temp_image.png", left, top, width=Inches(2), height=Inches(2))
                                            os.remove("temp_image.png")  # Clean up temp image
                                        else:
                                            continue
                                        run.text = new_text
            
            s_name = df.loc[i, event.email_column].split('@')[0]
            prs.save(os.path.join(root_directory, f"{s_name}.pptx"))
            f_id = ppt2pdf(os.path.join(root_directory, f"{s_name}.pptx"), s_name, token)
            download_url = f"https://docs.google.com/presentation/d/{f_id}/export/pdf"

            r = requests.get(download_url, allow_redirects=True)
            certificate_path = os.path.join(root_directory, f"{s_name}_certificate.pdf")
            open(certificate_path, 'wb').write(r.content)

            try:
                mail = EmailMessage(subject,
                                    f"Hello, {s_name} \n{mess}\n\nThis service is powered by EDEMS PVT Ltd.",
                                    settings.EMAIL_HOST_USER,
                                    [df.loc[i, event.email_column]])
                mail.attach_file(certificate_path)
                mail.send()
                Participant(event=event, email=df.loc[i, event.email_column], status=True).save()

                # Prepare response to download certificate
                response = FileResponse(open(certificate_path, 'rb'), content_type='application/pdf')
                response['Content-Disposition'] = f'attachment; filename="{s_name}_certificate.pdf"'

                # Remove certificate file after sending email
                os.remove(certificate_path)
                os.remove(os.path.join(root_directory, f"{s_name}.pptx"))

                return response  # Return response for download
            except Exception as e:
                Participant(event=event, email=df.loc[i, event.email_column], status=False).save()

            i += 1

        messages.success(request, "Certificates Sent Successfully !!")
        return redirect(f"/certificate/{event.id}/{event.slug}")

    return render(request, 'certificate/map_tags_of_template.html', {
        'columns': list(pd.read_csv(event.csv_file).columns),
        'tags': tags,
    })

@login_required
def view_certificate_status(request):
    return render(request, 'certificate/view_certificate_status.html', {
        'events': Event.objects.filter(user=request.user)
    })
